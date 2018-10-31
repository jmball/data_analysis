# Library of functions for analysis and plotting of the Automated PV
# Measurement data output.

import os
from functools import reduce

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib import gridspec
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from scipy.interpolate import interp1d
from scipy.signal import savgol_filter


def title_image_slide(prs, title):
    """
    Creates a new slide in the presentation (prs) with a formatted title.
    """

    # Add a title slide
    title_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_slide_layout)

    # Add text to title and edit its layout
    title_placeholder = slide.shapes.title
    title_placeholder.top = Inches(0)
    title_placeholder.width = Inches(10)
    title_placeholder.height = Inches(0.5)
    title_placeholder.text = title

    # Edit margins within textbox
    text_frame = title_placeholder.text_frame
    text_frame.margin_bottom = Inches(0)
    text_frame.margin_top = Inches(0.1)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Edit title fontsize and style
    p = text_frame.paragraphs[0]
    run = p.runs[0]
    font = run.font
    font.size = Pt(16)
    font.bold = True

    return slide


def extra_JV_analysis(filepath, intensity, num_cols):
    """Process JV curve data.

    Parameters
    ----------
    filepath : str
        absolute path to the data file
    intensity : float
        equivalent solar intensity (no. of suns)
    num_cols : int
        number of columns in log file

    Returns
    -------
    area : float
        area in cm^2
    scan_direction : str
        high-to-low ('HL') or low-to-high ('LH') scan direction
    condition : str
        'light' or 'dark' measurement condition
    jsc_int : float
        interpolated jsc in mA/cm^2
    voc_int : float
        interpolated voc in V
    ff_int : float
        fill factor derived from interpolated data
    pce_int : float
        power conversion efficiency derived from interpolated data in %
    vmp_int : float
        interpolated vmp in V
    jmp_int : float
        interpolated jmp in mA/cm^2
    vspo : float
        voltage used for spo scan
    rs_grad : float
        series resistance derived from gradient at voc in ohms
    rsh_grad : float
        shunt resistance derived from gradient at jsc in ohms
    rel_path : str
        relative filepath
    """

    # get relative file path
    rel_path = filepath.split('\\')[-1]
    # rel_path = filepath
    print(rel_path)

    # determine condition
    if rel_path.find('liv') > -1:
        condition = 'Light'
    elif rel_path.find('div') > -1:
        condition = 'Dark'
    elif rel_path.find('hold') > -1:
        condition = 'SPO'
    elif rel_path.find('jsc') > -1:
        condition = 'SJsc'
    elif rel_path.find('voc') > -1:
        condition = 'SVoc'

    # perform analysis depending on measurement condition
    if (condition == 'SPO') or (condition == 'SJsc') or (condition == 'SVoc'):
        # import datafile for retreiving metadata
        if (num_cols == 20) or (num_cols == 21):
            area = 0
            vspo = 0
        else:
            SPO = pd.read_csv(rel_path, delimiter='\t')
            area = SPO.iloc[-3, 1]
            vspo = SPO.iloc[-2, 1]
        scan_direction = '0'
        jsc_int = 0
        voc_int = 0
        ff_int = 0
        pce_int = 0
        vmp_int = 0
        jmp_int = 0
        rs_grad = 0
        rsh_grad = 0
    elif condition == 'Dark':
        # import J-V data and paramter metadata
        if (num_cols == 20) or (num_cols == 21):
            JV = np.genfromtxt(rel_path, delimiter='\t', skip_header=1, skip_footer=num_cols, usecols=(1, 3))
            area = 0
        else:
            JV = np.genfromtxt(rel_path, delimiter='\t')
            area = JV[-3, 1]
            intensity = JV[-1, 1]
            JV = JV[0:-11, :]

        # determine scan direction
        if JV[0, 0] > JV[-1, 0]:
            scan_direction = 'HL'
        else:
            scan_direction = 'LH'

        jsc_int = 0
        voc_int = 0
        ff_int = 0
        pce_int = 0
        vmp_int = 0
        jmp_int = 0
        vspo = 0
        rs_grad = 0
        rsh_grad = 0
    elif condition == 'Light':
        # import J-V data and paramter metadata
        if (num_cols == 20) or (num_cols == 21):
            JV = np.genfromtxt(rel_path, delimiter='\t', skip_header=1,
                                skip_footer=num_cols, usecols=(1, 2, 3, 4))
            area = 0
            I = JV[:,1]
            P = JV[:,3]
            JV = JV[:,:3:2]
        else:
            JV = np.genfromtxt(rel_path, delimiter='\t')
            area = JV[-3, 1]
            intensity = JV[-1, 1]
            JV = JV[0:-11, :]
            I = JV[:, 1] * area / 1000
            P = JV[:, 0] * JV[:, 1]

        # determine scan direction
        if JV[0, 0] > JV[-1, 0]:
            scan_direction = 'HL'
        else:
            scan_direction = 'LH'

        # calculate derived parameters for further analysis
        dpdv = np.gradient(P, JV[:, 0])
        didv = np.gradient(I, JV[:, 0])
        didv_sgfilter = savgol_filter(didv, 15, 3)
        R = 1 / didv_sgfilter

        # create interpolation functions
        try:
            f_j = interp1d(JV[:, 0], JV[:, 1], kind='cubic', bounds_error=False, fill_value=0)
            jsc_int = np.float64(f_j(0))
        except ValueError:
            jsc_int = 0
        try:
            f_v = interp1d(JV[:, 1], JV[:, 0], kind='cubic', bounds_error=False, fill_value=0)
            voc_int = np.float64(f_v(0))
        except ValueError:
            voc_int = 0
        try:
            f_p = interp1d(JV[:, 0], P, kind='cubic', bounds_error=False, fill_value=0)
            f_dpdv = interp1d(dpdv, JV[:, 0], kind='cubic', bounds_error=False, fill_value=0)
            vmp_int = np.float64(f_dpdv(0))
            jmp_int = np.float64(f_j(vmp_int))
            ff_int = jmp_int * vmp_int / (jsc_int * voc_int)
            pmax = np.float64(f_p(vmp_int))
        except ValueError:
            vmp_int = 0
            jmp_int = 0
            ff_int = 0
            pmax = 0
        try:
            f_r = interp1d(JV[:, 0], R, kind='cubic', bounds_error=False, fill_value=0)
            rsh_grad = np.absolute(np.float64(f_r(0)))
            rs_grad = np.absolute(np.float64(f_r(voc_int)))
        except ValueError:
            rsh_grad = 0
            rs_grad = 0
        pce_int = pmax / float(intensity)
        vspo = 0

    return [
        area, scan_direction, condition, jsc_int, voc_int, ff_int, pce_int,
        vmp_int, jmp_int, vspo, rs_grad, rsh_grad, rel_path
    ]


def boxplotdata(grouped_by_var):
    """

    Function for sorting log file data grouped by variable into series for
    boxplots.

    """

    Jsc_var = []
    Voc_var = []
    FF_var = []
    PCE_var = []
    Rs_var = []
    Rsh_var = []
    names_var = []
    var_names = []
    for name, group in grouped_by_var:
        grouped_by_val = group.groupby('Value')
        var_names.append(name)
        Jsc = []
        Voc = []
        FF = []
        PCE = []
        Rs = []
        Rsh = []
        names_val = []
        for name, group in grouped_by_val:
            names_val.append(name)
            Jsc.append(np.array(group['Jsc_int']))
            Voc.append(np.array(group['Voc_int']))
            FF.append(np.array(group['FF_int']))
            PCE.append(np.array(group['PCE_int']))
            Rs.append(np.array(group['Rs_grad']))
            Rsh.append(np.array(group['Rsh_grad']))
        Jsc_var.append(Jsc)
        Voc_var.append(Voc)
        FF_var.append(FF)
        PCE_var.append(PCE)
        Rs_var.append(Rs)
        Rsh_var.append(Rsh)
        names_var.append(names_val)
    return {
        'Jsc_var': Jsc_var,
        'Voc_var': Voc_var,
        'FF_var': FF_var,
        'PCE_var': PCE_var,
        'Rs_var': Rs_var,
        'Rsh_var': Rsh_var,
        'names_var': names_var,
        'var_names': var_names
    }


def create_figure(rows,
                  cols,
                  title=None,
                  width=4.75,
                  height=3.5625,
                  dpi=300,
                  title_fsize=10,
                  title_fweight='bold'):
    """
    Create a figure object (fig), give it a title, and create a gridspec obejct
    (gs). Subplots are created by fig in positions limited by the dimentions of
    gs.

    title = title of the full page figure
    rows = number of subplot rows
    cols = number of subplot columns
    width = figure width (default 10 is the width of an A4 page in inches)
    height = figure height (default 7.5 is the height of an A4 page in inches)
    dpi = dots per inches of figure output
    title_fsize = fontsize of the title
    """

    fig = plt.figure(figsize=(width, height), dpi=dpi)
    if title is not None:
        fig.suptitle(title, fontsize=title_fsize, fontweight=title_fweight)
    gs = gridspec.GridSpec(rows, cols)
    return fig, gs


def subboxplot(self,
               boxplotdata_HL,
               boxplotdata_LH,
               scatter_x,
               label_HL,
               label_LH,
               set_yscale='linear',
               showfliers=False,
               boxprops_HL=dict(color='blue'),
               whiskerprops_HL=dict(color='blue', linestyle='-'),
               capprops_HL=dict(color='blue'),
               medianprops_HL=dict(color='blue'),
               boxprops_LH=dict(color='red'),
               whiskerprops_LH=dict(color='red', linestyle='-'),
               capprops_LH=dict(color='red'),
               medianprops_LH=dict(color='red')):
    """
    Method for creating a formatted subplot boxplot.

    This function will be bound to the matplotlib.axes.SubplotBase class as
    a new method for instances of the class. It should only be called when
    binding the method and after only as a method of SubplotBase objects.
    """

    self.boxplot(
        boxplotdata_HL,
        showfliers=showfliers,
        labels=label_HL,
        boxprops=boxprops_HL,
        whiskerprops=whiskerprops_HL,
        capprops=capprops_HL,
        medianprops=medianprops_HL)
    self.boxplot(
        boxplotdata_LH,
        showfliers=showfliers,
        labels=label_LH,
        boxprops=boxprops_LH,
        whiskerprops=whiskerprops_LH,
        capprops=capprops_LH,
        medianprops=medianprops_LH)
    self.set_xticklabels(label_HL, rotation=45, ha='right')
    self.scatter(
        scatter_x,
        np.concatenate(boxplotdata_HL),
        c='blue',
        marker='x',
        label='H->L')
    self.scatter(
        scatter_x,
        np.concatenate(boxplotdata_LH),
        c='red',
        marker='x',
        label='L->H')
    self.set_yscale(set_yscale)


def subbarchart(self, x, y, names):
    """
    Method for creating a formatted subplot bar chart.

    This function will be bound to the matplotlib.axes.SubplotBase class as
    a new method for instances of the class. It should only be called when
    binding the method and after only as a method of SubplotBase objects.
    """

    self.bar(
        x, y, align='center', width=0.4, edgecolor='black', color='black')
    self.set_xticks(x)
    self.set_xticklabels(names, fontdict={'fontsize': 'xx-small'}, rotation=45, ha='right')


def set_axes_props(self,
                   title_fsize=10,
                   title=None,
                   xlabel=None,
                   xscale=None,
                   ylabel=None,
                   yscale=None,
                   xlim=None,
                   ylim=None,
                   axhline=None,
                   axvline=None,
                   xticks=None,
                   yticks=None,
                   xticklabels=None,
                   yticklabels=None,
                   legend=False,
                   loc='best',
                   scatterpoints=1,
                   fontsize=8):
    """
    Function for setting several properties of the axes of a matplotlib subplot
    at once.

    See http://matplotlib.org/api/axes_api.html for further details of each
    option.

    This function will be bound to the matplotlib.axes.SubplotBase class as
    a new method for instances of the class. It should only be called when
    binding the method and then only as a method of SubplotBase objects.
    """

    if xlabel is not None:
        self.set_xlabel(xlabel)
    if ylabel is not None:
        self.set_ylabel(ylabel)
    if xscale is not None:
        self.set_xscale(xscale)
    if yscale is not None:
        self.set_yscale(yscale)
    if xlim is not None:
        self.set_xlim(xlim)
    if ylim is not None:
        self.set_ylim(ylim)
    if axhline is not None:
        self.axhline(axhline, lw=0.5, c='black')
    if axvline is not None:
        self.axvline(axvline, lw=0.5, c='black')
    if xticks is not None:
        self.set_xticks(xticks)
    if xticklabels is not None:
        self.set_xticklabels(xticklabels)
    if yticks is not None:
        self.set_yticks(yticks)
    if yticklabels is not None:
        self.set_yticklabels(yticklabels)
    if legend is True:
        self.legend(loc=loc, fontsize=fontsize, scatterpoints=scatterpoints)
    if title is not None:
        self.set_title(title, fontsize=title_fsize)


def save_image(gs, fig, image_path, wspace, hspace):
    """
    Save an image of fig to a file.
    """

    # plt.subplots_adjust(top=0.92)
    gs.update(wspace=wspace, hspace=hspace)
    fig.tight_layout()
    fig.savefig(image_path)


def create_save_boxplot(parameter,
                        boxplotdata_HL,
                        boxplotdata_LH,
                        index,
                        scatter_x,
                        image_path,
                        title=None,
                        showfliers=False,
                        boxprops_HL=dict(color='black', linewidth=0.5),
                        whiskerprops_HL=dict(color='black', linestyle='-', linewidth=0.5),
                        capprops_HL=dict(color='black', linewidth=0.5),
                        medianprops_HL=dict(color='black', linewidth=0.5),
                        boxprops_LH=dict(color='red', linewidth=0.5),
                        whiskerprops_LH=dict(color='red', linestyle='-', linewidth=0.5),
                        capprops_LH=dict(color='red', linewidth=0.5),
                        medianprops_LH=dict(color='red', linewidth=0.5)):
    """
    Create a figure object (fig), add boxplots for each variable, format
    the axes depending on the PV parameter, and save it to a file.
    """
    fig = plt.figure(figsize=(4.75, 3.5625), dpi=300)
    ax = fig.add_subplot(1, 1, 1)
    p = parameter + '_var'

    ax.boxplot(
        np.absolute(boxplotdata_HL[p][index]),
        showfliers=showfliers,
        labels=boxplotdata_HL['names_var'][index],
        boxprops=boxprops_HL,
        whiskerprops=whiskerprops_HL,
        capprops=capprops_HL,
        medianprops=medianprops_HL)
    ax.boxplot(
        np.absolute(boxplotdata_LH[p][index]),
        showfliers=showfliers,
        labels=boxplotdata_LH['names_var'][index],
        boxprops=boxprops_LH,
        whiskerprops=whiskerprops_LH,
        capprops=capprops_LH,
        medianprops=medianprops_LH)
    ax.set_xticklabels(
        boxplotdata_HL['names_var'][index], fontdict={'fontsize': 'xx-small'}, rotation=45, ha='right')
    ax.scatter(
        scatter_x,
        np.absolute(np.concatenate(boxplotdata_HL[p][index])),
        s=3,
        c='black',
        marker='o',
        label='H->L')
    ax.scatter(
        scatter_x,
        np.absolute(np.concatenate(boxplotdata_LH[p][index])),
        s=3,
        c='red',
        marker='o',
        label='L->H')
    if parameter == 'Jsc':
        ax.set_ylabel('Jsc (mA/cm^2)')
        ax.set_ylim(0)
    elif parameter == 'Voc':
        ax.set_ylabel('Voc (V)')
        ax.set_ylim(0)
    elif parameter == 'FF':
        ax.set_ylabel('FF')
        ax.set_ylim((0, 1))
    elif parameter == 'PCE':
        ax.set_ylabel('PCE (%)')
        ax.set_ylim(0)
    elif parameter == 'Rs':
        ax.set_ylabel('Rs (ohms)')
        ax.set_yscale('log')
    elif parameter == 'Rsh':
        ax.set_ylabel('Rsh (ohms)')
        ax.set_yscale('log')
    ax.legend(loc='best', fontsize='xx-small', scatterpoints=1)
    fig.tight_layout()
    fig.savefig(image_path)

    # np.savetxt(
    #     image_path.strip('.png') + '_HL.txt',
    #     np.transpose(boxplotdata_HL[p][index]),
    #     header='\t'.join(boxplotdata_HL['names_var'][index]),
    #     comments='',
    #     newline='\r\n',
    #     delimiter='\t')


def create_save_barchart(yields, names, index, image_path, title=None):
    """
    Create a figure object (fig), add a bar chart, format the axes, and save
    it to a file.
    """

    fig = plt.figure(figsize=(4.75, 3.5625), dpi=300)
    ax = fig.add_subplot(1, 1, 1)
    ax.bar(
        range(len(yields[index])),
        yields[index],
        align='center',
        width=0.4,
        edgecolor='black',
        color='black')
    ax.set_xticks(range(len(yields[index])))
    ax.set_xticklabels(names[index], fontdict={'fontsize': 'xx-small'}, rotation=45, ha='right')
    ax.set_ylim([0, 105])
    ax.set_ylabel('Yield (%)')
    fig.tight_layout()
    fig.savefig(image_path)


def exp_file_list(folderpath):
    """

    Determine whether data exists for an additional experiment. If so, make
    a list of files. Returns a dictionary containing a boolean indicating
    existance and a list of files (empty if experiment hasn't been performed).

    """

    files = []
    files_split = []
    if os.path.exists(folderpath):
        exists = True
        for f in os.listdir(folderpath):
            if (f.endswith('.txt')) & (f.find('_LOG.txt') == -1):
                files.append(folderpath + f)
                files_split.append(f.split('_'))
        if len(files) == 0:
            exists = False
    else:
        exists = False

    return {'exists': exists, 'files': files, 'files_split': files_split}


def build_log_df(exp_files, master_log):
    """

    Build a log dataframe from from experiment files and master J-V log.

    """

    # Make a dataframe containing labels, pixels, and scan directions based on
    # the filename of each data file.
    labels = []
    pixels = []
    scan_dir = []
    for item in exp_files['files_split']:
        labels.append(item[0])
        pixels.append(item[1])
        scan_dir.append(item[3])
    files_dict = {
        'Label': labels,
        'Pixel': pixels,
        'File_Path': exp_files['files'],
        'Scan_direction': scan_dir
    }
    files_df = pd.DataFrame(files_dict)
    files_df = files_df.sort_values(['Label'], ascending=True)

    # Find all of the unique labels in the master JV log file. Convert 'Label'
    # series to strings to ensure the types match when checking if a label
    # also appears in files_df.
    unique_labels_df = master_log.drop_duplicates(['Label'])
    unique_labels_df = unique_labels_df.sort_values(['Label'], ascending=True)
    unique_labels_df['Label'] = unique_labels_df['Label'].astype('str')
    unique_labels_df = unique_labels_df[unique_labels_df.Label.isin(
        files_df.Label.values)]

    # Reset index of unique_labels_df to ensure they are added in the correct
    # row of files_df.
    unique_labels_df.reset_index(drop=True, inplace=True)
    files_df['Variable'] = unique_labels_df['Variable']
    files_df['Value'] = unique_labels_df['Value']

    return files_df


def factors(n):
    """

    Returns a set of factors of the integer n.

    """

    pairs = [[i, n // i] for i in range(1, int(n**0.5) + 1) if n % i == 0]
    return set(reduce(list.__add__, pairs))


def round_sig_fig(x, sf):
    """

    Rounds any number, x, to the specified number of significant figures, sf.

    """

    format_str = '%.' + str(sf) + 'e'
    x_dig, x_ord = map(float, (format_str % x).split('e'))
    return round(x, int(-x_ord) + 1)


def calc_ticks(data_min, data_max):
    """

    Returns an array of axis tick labels given the max and min values of
    the data to be plotted.

    """

    min_digs, min_order = map(float, ('%.6e' % data_min).split('e'))
    max_digs, max_order = map(float, ('%.6e' % data_max).split('e'))
    min_order -= 1
    max_order -= 1
    if max_order > min_order:
        min_digs = np.floor(min_digs * 10**(min_order - max_order + 1))
        max_digs = np.ceil(max_digs * 10)
        y_range_digs = max_digs - min_digs
        while np.max(list(filter(lambda x: x < 8, factors(y_range_digs)))) < 3:
            y_range_digs += 2
            min_digs -= 1
            max_digs += 1
        num_ticks = np.max(
            list(filter(lambda x: x < 8, factors(y_range_digs))))
        min_tick = min_digs * 10**max_order
        ticks = np.linspace(min_tick, y_range_digs * 10**max_order + min_tick,
                            num_ticks + 1)
    elif min_order > max_order:
        min_digs = np.floor(min_digs * 10)
        max_digs = np.ceil(max_digs * 10**(max_order - min_order + 1))
        y_range_digs = max_digs - min_digs
        while np.max(list(filter(lambda x: x < 8, factors(y_range_digs)))) < 3:
            y_range_digs += 2
            min_digs -= 1
            max_digs += 1
        num_ticks = np.max(
            list(filter(lambda x: x < 8, factors(y_range_digs))))
        min_tick = min_digs * 10**min_order
        ticks = np.linspace(min_tick, y_range_digs * 10**min_order + min_tick,
                            num_ticks + 1)
    elif max_order == min_order:
        str_min = str(min_digs).replace('.', '').replace('-', '')
        str_max = str(max_digs).replace('.', '').replace('-', '')
        i = 0
        while str_min[i] == str_max[i]:
            i += 1
            if i == len(str_min):
                break
        if i == len(str_min):
            ticks = np.array([min_digs * 10**(min_order + 1)])
        else:
            min_digs = np.floor(min_digs * 10**(1 + i))
            max_digs = np.ceil(max_digs * 10**(1 + i))
            y_range_digs = max_digs - min_digs
            while np.max(list(filter(lambda x: x < 8,
                                     factors(y_range_digs)))) < 3:
                y_range_digs += 2
                min_digs -= 1
                max_digs += 1
            num_ticks = np.max(
                list(filter(lambda x: x < 8, factors(y_range_digs))))
            min_tick = min_digs * 10**(min_order - i)
            ticks = np.linspace(min_tick, y_range_digs * 10**
                                (min_order - i) + min_tick, num_ticks + 1)

    return ticks
