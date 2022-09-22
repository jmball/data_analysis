"""Generate a report of solar cell measurement data."""

import contextlib
import logging
import os
import pathlib
import time
import warnings

import matplotlib

matplotlib.use("TkAgg")

import matplotlib.pyplot as plt
import numpy as np
import packaging.version
import pandas as pd
import seaborn as sns
import scipy.constants
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR

from gooey import Gooey, GooeyParser

from check_release_version import get_latest_release_version, REPO_URL
from format_data import format_folder
from log_generator import generate_log
from version import __version__


# supress warnings
warnings.filterwarnings("ignore")

# Define a colormap for graded plots
cmap = plt.cm.get_cmap("viridis")


@Gooey(
    dump_build_config=False,
    program_name="Data Analysis",
    default_size=(750, 530),
    header_bg_color="#7B7B7B",
)
def parse():
    """Parse command line arguments to Gooey GUI."""
    desc = "Analyse solar simulator data and generate a report."

    # check if latest release on github is newer than currently running version
    # if so, let the user know by editing the description string
    latest_release_version = get_latest_release_version()
    if latest_release_version is None:
        desc += (
            "\n\nCould not determine latest release version. Check internet connection."
        )
    elif packaging.version.parse(latest_release_version) > packaging.version.parse(
        __version__
    ):
        desc += f"\n\nNEW VERSION AVAILABLE! Download it from: {REPO_URL}"
    else:
        desc += f"\n\nYou're running the latest version: {__version__}"

    parser = GooeyParser(description=desc)
    req = parser.add_argument_group(gooey_options={"columns": 1})
    req.add_argument(
        "folder",
        metavar="Folder containing data to be analysed",
        help="Absolute path to the folder containing measurement data",
        widget="DirChooser",
    )
    req.add_argument(
        "fix_ymin_0",
        metavar="Zero y-axis minima",
        help="Fix boxplot y-axis minima to 0",
        widget="Dropdown",
        choices=["yes", "no"],
        default="yes",
    )
    req.add_argument(
        "--debug",
        metavar="DEBUG",
        help="Export debug info to a file",
        widget="CheckBox",
        action="store_true",
    )
    return parser.parse_args()


def create_logger(log_dir: str, debug: bool = False):
    """Create a logger.

    Parameters
    ----------
    log_dir : str
        Log directory.
    debug : bool
        Flag whether to log in debug mode, which exports logging to a file.
    """
    # create logger
    logging.captureWarnings(True)
    _logger = logging.getLogger()
    log_level = 10 if debug else 20
    _logger.setLevel(log_level)

    # create a filter to remove messages from certain imports
    class ImportFilter(logging.Filter):
        """Filter log records from named third-party imports."""

        def filter(self, record: logging.LogRecord) -> bool:
            if record.name.startswith("matplotlib"):
                return False
            elif record.name.startswith("PIL"):
                return False
            else:
                return True

    # create console handler
    console_handler = logging.StreamHandler()
    console_handler.setLevel(log_level)
    console_handler.addFilter(ImportFilter())
    _logger.addHandler(console_handler)

    # add file handler for debugging
    if debug:
        formatter = logging.Formatter("%(asctime)s|%(name)s|%(levelname)s|%(message)s")

        file_handler = logging.FileHandler(
            pathlib.Path(log_dir).joinpath("debug.txt"), mode="w"
        )
        file_handler.setLevel(log_level)
        file_handler.setFormatter(formatter)
        file_handler.addFilter(ImportFilter())
        _logger.addHandler(file_handler)

    return _logger


def round_sig_fig(number: float, sig_fig: int) -> float:
    """
    Round a number to the specified number of significant figures.

    Parameters
    ----------
    number : float
        number to round
    sig_fig : int
        number of significant figures

    Returns
    -------
    rounded_number : float
        rounded number
    """
    _, x_ord = map(float, f"{number:.{sig_fig}e}".split("e"))
    return round(number, int(-x_ord) + 1)


def recursive_path_split(filepath) -> tuple:
    """Recursively split filepath into sub-parts delimited by OS file seperator.

    Parameters
    ----------
    path : str
        filepath

    Returns
    -------
    split : tuple
        sub-parts of filepath
    """
    head, tail = os.path.split(filepath)
    return (head,) if tail == "" else recursive_path_split(head) + (tail,)


def title_image_slide(prs, title: str):
    """
    Create a new slide in the presentation (prs) with a formatted title.

    Parameters
    ----------
    prs : presentation object
        pptx presentation object
    title : str
        title of slide

    Returns
    -------
    slide : slide object
        pptx slide object
    """
    # Add a title slide
    title_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_slide_layout)

    # Add text to title and edit its layout
    title_placeholder = slide.shapes.title
    title_placeholder.top = Inches(0)
    title_placeholder.width = Inches(10)
    title_placeholder.height = Inches(0.5)
    title_placeholder.text = title

    # Edit margins within textbox
    text_frame = title_placeholder.text_frame
    text_frame.margin_bottom = Inches(0)
    text_frame.margin_top = Inches(0.1)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Edit title fontsize and style
    para = text_frame.paragraphs[0]
    run = para.runs[0]
    font = run.font
    font.size = Pt(16)
    font.bold = True

    return slide


def plot_boxplots(
    dataframe,
    params,
    kind: str,
    grouping: str,
    variable: str = "",
    start_index: int = 0,
    data_slide=None,
    override_grouping_title: str = "",
):
    """Create boxplots from the log file.

    Parameters
    ----------
    dataframe : pandas.DataFrame
        logfile or group
    params : list of str
        parameters to plot
    kind : str
        kind of paramters plotted
    grouping : str
        how data are grouped
    variable : str
        variable
    start_index : int
        starting index of boxplot. Useful if carrying on page from previous plots.
    data_slide : prs object
        current slide. Useful if carrying on from previous plots.
    override_grouping_title : str
        instead of using `grouping` in slide title, use this string.
    """
    plot_labels_dict = {
        "jsc": {"J-V": "Jsc (mA/cm^2)"},
        "voc": {"J-V": "Voc (V)"},
        "pce": {"J-V": "PCE (%)"},
        "vmpp": {"J-V": "Vmp (V)"},
        "jmpp": {"J-V": "Jmp (mA/cm^2)"},
        "jss": {"SSPO": "J_mp_ss (mA/cm^2)", "SSJsc": "J_sc_ss (mA/cm^2)"},
        "pcess": {"SSPO": "PCE_ss (%)"},
        "vss": {"SSPO": "V_mp_ss (V)", "SSVoc": "V_oc_ss (V)"},
        "ff": {"J-V": "FF"},
        "quasiff": {"SSJsc": "Quasi-FF"},
        "rsvfwd": {"J-V": "Rs (ohms)"},
        "rsh": {"J-V": "Rsh (ohms)"},
        "pcesspcejv": {"SSPO": "PCE_ss/PCE_jv"},
    }

    plot_index = 0
    for param in params:
        # create a new slide for every 4 plots
        if (start_index + plot_index) % 4 == 0:
            ss_or_jv = kind if kind == "J-V" else "Steady-state"
            grouping_title = override_grouping_title or grouping
            page = int((start_index + plot_index) / 4)
            data_slide = title_image_slide(
                prs,
                f"{variable} {ss_or_jv} parameters by {grouping_title}, page {page}",
            )

        # create boxplot
        fig, ax1 = plt.subplots(
            1, 1, dpi=300, **{"figsize": (A4_WIDTH / 2, A4_HEIGHT / 2)}
        )

        # get grouping of data for box and swarm plots
        if kind == "J-V":
            hue = (
                dataframe["scandirection"]
                + np.array([", "] * len(dataframe["area"]))
                + dataframe["area"].astype(str)
            )
        else:
            hue = dataframe["area"]

        try:
            sns.boxplot(
                x=dataframe[grouping],
                y=np.absolute(dataframe[param].astype(float)),
                hue=hue,
                palette="deep",
                linewidth=0.5,
                ax=ax1,
                showfliers=False,
            )
        except ValueError as err:
            logger.error(hue, grouping, param, kind)
            logger.error(dataframe["jsc"])
            raise ValueError from err
        sns.swarmplot(
            x=dataframe[grouping],
            y=np.absolute(dataframe[param].astype(float)),
            hue=hue,
            palette="muted",
            size=3,
            linewidth=0.5,
            edgecolor="gray",
            dodge=True,
            ax=ax1,
        )

        # only show legend markers for box plots, not swarm plot
        legend_handles, legend_labels = ax1.get_legend_handles_labels()
        ax1.legend(
            legend_handles[: len(legend_handles) // 2],
            legend_labels[: len(legend_labels) // 2],
            fontsize="small",
        )

        ax1.set_xticklabels(
            ax1.get_xticklabels(), fontsize="small", rotation=45, ha="right"
        )
        ax1.set_xlabel("")
        if param in ["jsc", "voc", "pce", "vmpp", "jmpp", "jss", "pcess", "vss"]:
            if FIX_YMIN_0:
                ax1.set_ylim(0)
        elif param in ["ff", "quasiff"]:
            if FIX_YMIN_0:
                ax1.set_ylim(0, 1)

        ax1.set_ylabel(plot_labels_dict[param][kind], fontsize="small")

        fig.tight_layout()

        # save figure and add to powerpoint
        image_png = os.path.join(image_folder, f"boxplot_{param}.png")
        image_svg = os.path.join(image_folder, f"boxplot_{param}.svg")
        fig.savefig(image_png)
        fig.savefig(image_svg)
        if data_slide is not None:
            data_slide.shapes.add_picture(
                image_png,
                left=LEFTS[str((start_index + plot_index) % 4)],
                top=TOPS[str((start_index + plot_index) % 4)],
                height=IMAGE_HEIGHT,
            )
        plot_index += 1

    return start_index + plot_index, data_slide


def plot_countplots(
    dataframe, index: int, grouping: str, data_slide, variable: str = ""
):
    """Create countplots from the log file.

    Parameters
    ----------
    dataframe : DataFrame
        logfile or group
    index : int
        figure index
    grouping : str
        how data are grouped
    data_slide: slide
        slide in ppt to add figures to
    variable : str
        variable
    """
    # create count plot
    fig, ax1 = plt.subplots(1, 1, dpi=300, **{"figsize": (A4_WIDTH / 2, A4_HEIGHT / 2)})
    if grouping == "value":
        ax1.set_title(f"{variable}", fontdict={"fontsize": "small"})
    sns.countplot(
        x=dataframe[grouping],
        data=dataframe,
        hue=dataframe["scandirection"],
        linewidth=0.5,
        palette="deep",
        edgecolor="black",
        ax=ax1,
    )
    legend_handles, legend_labels = ax1.get_legend_handles_labels()
    ax1.legend(
        legend_handles[:],
        legend_labels[:],
        fontsize="small",
    )
    ax1.set_xticklabels(
        ax1.get_xticklabels(), fontsize="small", rotation=45, ha="right"
    )
    ax1.set_xlabel("")
    ax1.set_ylabel("Number of working pixels", fontsize="small")
    fig.tight_layout()

    # save figure and add to powerpoint
    image_png = os.path.join(image_folder, f"boxchart_yields{index}.png")
    image_svg = os.path.join(image_folder, f"boxchart_yields{index}.svg")
    fig.savefig(image_png)
    fig.savefig(image_svg)
    data_slide.shapes.add_picture(
        image_png,
        left=LEFTS[str(index % 4)],
        top=TOPS[str(index % 4)],
        height=IMAGE_HEIGHT,
    )


def plot_stabilisation(dataframe, title: str, short_name: str):
    """Plot stabilisation data.

    Parameters
    ----------
    dataframe : dataFrame
        data to plot
    title : str
        slide title
    short_name : str
        short name for file
    """
    for index, (_, row) in enumerate(dataframe.iterrows()):
        # Get label, variable, value, and pixel for title and image path
        label = row["label"]
        variable = row["variable"]
        value = row["value"]
        pixel = row["pixel"]
        vspo = row["vss"]

        # Start a new slide after every 4th figure
        if index % 4 == 0:
            data_slide = title_image_slide(prs, f"{title}, page {int(index / 4)}")

        # Open the data file
        path = row["relativepath"]
        if short_name == "sjsc":
            cols = (0, 4)
        elif short_name == "spo":
            cols = (0, 2, 4, 6)
        elif short_name == "svoc":
            cols = (0, 2)
        else:
            cols = ()

        data = np.genfromtxt(
            path, delimiter="\t", skip_header=1, skip_footer=NUM_COLS, usecols=cols
        )
        with contextlib.suppress(Exception):
            data = data[~np.isnan(data).any(axis=1)]

        try:
            if short_name == "sjsc":
                fig, ax1 = plt.subplots(
                    1, 1, figsize=(A4_WIDTH / 2, A4_HEIGHT / 2), dpi=300
                )
                ax1.set_title(
                    f"{label}, pixel {pixel}, {variable}, {value}",
                    fontdict={"fontsize": "small"},
                )
                ax1.scatter(
                    data[:, 0], np.absolute(data[:, 1]), color="black", s=5, label="Jsc"
                )
                ax1.set_ylabel("|Jsc| (mA/cm^2)", fontsize="small")
                ax1.set_ylim(0, np.max(np.absolute(data[:, 1])) * 1.1)
                ax1.set_xlabel("Time (s)", fontsize="small")
                ax1.set_xlim(0)
                ax1.tick_params(direction="in", top=True, right=True, labelsize="small")
                fig.tight_layout()
            elif short_name == "spo":
                fig, axs = plt.subplots(
                    3, 1, sharex=True, figsize=(A4_WIDTH / 2, A4_HEIGHT / 2), dpi=300
                )
                ax1, ax2, ax3 = axs
                fig.subplots_adjust(hspace=0)
                ax1.set_title(
                    f"{label}, pixel {pixel}, {variable}, {value}, vspo = {vspo} V",
                    fontdict={"fontsize": "small"},
                )
                ax1.scatter(
                    data[:, 0], np.absolute(data[:, 2]), color="black", s=5, label="J"
                )
                ax1.set_ylabel("|J| (mA/cm^2)", fontsize="small")
                ax1.set_ylim(0, np.max(np.absolute(data[:, 2])) * 1.1)
                ax3.tick_params(direction="in", top=True, right=True)
                ax2.scatter(
                    data[:, 0],
                    np.absolute(data[:, 3]),
                    color="red",
                    s=5,
                    marker="s",
                    label="pce",
                )
                ax2.set_ylabel("PCE (%)", fontsize="small")
                ax2.set_ylim(0, np.max(np.absolute(data[:, 3])) * 1.1)
                ax2.tick_params(direction="in", top=True, right=True)
                ax3.scatter(
                    data[:, 0],
                    np.absolute(data[:, 1]),
                    color="blue",
                    s=5,
                    marker="s",
                    label="v",
                )
                ax3.set_ylabel("V (V)", fontsize="small")
                ax3.set_ylim(0, np.max(np.absolute(data[:, 1])) * 1.1)
                ax3.set_xlabel("Time (s)", fontsize="small")
                ax3.tick_params(direction="in", top=True, right=True, labelsize="small")
                fig.align_ylabels([ax1, ax2, ax3])
            elif short_name == "svoc":
                fig, ax1 = plt.subplots(
                    1, 1, figsize=(A4_WIDTH / 2, A4_HEIGHT / 2), dpi=300
                )
                ax1.set_title(
                    f"{label}, pixel {pixel}, {variable}, {value}",
                    fontdict={"fontsize": "small"},
                )
                ax1.scatter(
                    data[:, 0], np.absolute(data[:, 1]), color="black", s=5, label="Voc"
                )
                ax1.set_ylabel("|Voc| (V)", fontsize="small")
                ax1.set_ylim(0, np.max(np.absolute(data[:, 1])) * 1.1)
                ax1.set_xlabel("Time (s)", fontsize="small")
                ax1.set_xlim(0)
                ax1.tick_params(direction="in", top=True, right=True, labelsize="small")
                fig.tight_layout()
            else:
                fig, ax1 = plt.subplots(
                    1, 1, figsize=(A4_WIDTH / 2, A4_HEIGHT / 2), dpi=300
                )

            # Format the figure layout, save to file, and add to ppt
            image_png = os.path.join(
                image_folder, f"{short_name}_{label}_{variable}_{value}_{pixel}.png"
            )
            image_svg = os.path.join(
                image_folder, f"{short_name}_{label}_{variable}_{value}_{pixel}.svg"
            )
            fig.savefig(image_png)
            fig.savefig(image_svg)
            data_slide.shapes.add_picture(
                image_png,
                left=LEFTS[str(index % 4)],
                top=TOPS[str(index % 4)],
                height=IMAGE_HEIGHT,
            )

            # Close figure
            plt.close(fig)
        except IndexError:
            logger.error("indexerror")

        index += 1


def plot_spectra(files):
    """Plot illumination specta.

    Parameters
    ----------
    files : list
        list of file paths
    """
    c_div = 1 / len(files)

    data_slide = title_image_slide(prs, "Measured illumination spectra")

    fig, ax1 = plt.subplots(1, 1, figsize=(A4_WIDTH, A4_HEIGHT), dpi=300)
    for index, file in enumerate(files):
        if os.path.getsize(file) != 0:
            spectrum = np.genfromtxt(file, delimiter="\t")
            ax1.plot(
                spectrum[:, 0],
                spectrum[:, 1],
                color=cmap(index * c_div),
                label=f"{index}",
            )
    ax1.set_ylabel("Spectral irradiance (W/cm^2/nm)", fontsize="large")
    ax1.set_ylim(0)
    ax1.tick_params(direction="in", top=True, right=True, labelsize="large")
    ax1.set_xlabel("Wavelength (nm)", fontsize="large")
    ax1.set_xlim(350, 1100)
    ax1.legend(fontsize="large")

    fig.tight_layout()

    # Format the figure layout, save to file, and add to ppt
    image_png = os.path.join(image_folder, "spectra.png")
    image_svg = os.path.join(image_folder, "spectra.svg")
    fig.savefig(image_png)
    fig.savefig(image_svg)
    data_slide.shapes.add_picture(
        image_png, left=LEFTS["0"], top=TOPS["0"], height=IMAGE_HEIGHT * 2
    )

    # Close figure
    plt.close(fig)


def plot_best_jvs_by_label(groups, substrate_info):
    """Plot best jv curves for each substrate.

    Parameters
    ----------
    groups : pandas.GroupBy
        data frame grouped by label.
    substrate_info : pandas.DataFrame
        data frame from which variables, values, and labels can be inferred.
    """
    variables = list(substrate_info["variable"])
    values = list(substrate_info["value"])
    labels = list(substrate_info["label"])

    # get parameters for plot formatting
    c_div = 1 / 8

    # Create figures, save images and add them to powerpoint slide
    for index, (_, group) in enumerate(groups):
        # Create a new slide after every four graphs are produced
        if index % 4 == 0:
            data_slide = title_image_slide(
                prs, f"Best JV scans of every working pixel, page {int(index / 4)}"
            )

        # Create figure, axes, y=0 line, and title
        fig, ax1 = plt.subplots(1, 1, figsize=(A4_WIDTH / 2, A4_HEIGHT / 2), dpi=300)
        ax1.axhline(0, lw=0.5, c="black")
        ax1.axvline(0, lw=0.5, c="black")
        ax1.set_title(
            f"{labels[index]}, {variables[index]}, {values[index]}",
            fontdict={"fontsize": "small"},
        )

        pixels = list(group["pixel"].astype(int))

        # find signs of jsc and voc to determine max and min axis limits
        jsc_signs, jsc_counts = np.unique(np.sign(group["jmpp"]), return_counts=True)
        voc_signs, voc_counts = np.unique(np.sign(group["voc"]), return_counts=True)
        if len(jsc_signs) == 1:
            jsc_sign = jsc_signs[0]
        else:
            max_ix = np.argmax(jsc_counts)
            jsc_sign = jsc_signs[max_ix]
        if len(voc_signs) == 1:
            voc_sign = voc_signs[0]
        else:
            max_ix = np.argmax(voc_counts)
            voc_sign = voc_signs[max_ix]

        # load data for each pixel and plot on axes
        fwd_j = []
        rev_j = []
        for data_index, (file, scan_dir) in enumerate(
            zip(group["relativepath"], group["scandirection"])
        ):
            if scan_dir == "fwd":
                data_fwd = np.genfromtxt(
                    file,
                    delimiter="\t",
                    skip_header=1,
                    skip_footer=NUM_COLS,
                    usecols=(2, 4),
                )
                data_fwd = data_fwd[~np.isnan(data_fwd).any(axis=1)]
                ax1.plot(
                    data_fwd[:, 0],
                    data_fwd[:, 1],
                    c=cmap(pixels[data_index] * c_div),
                    lw=2.0,
                )
                if (
                    (jsc_sign > 0) & (voc_sign > 0)
                    or not (jsc_sign > 0) & (voc_sign < 0)
                    and (jsc_sign < 0) & (voc_sign > 0)
                ):
                    fwd_j.append(data_fwd[0, 1])
                    rev_j.append(data_fwd[-1, 1])
                elif (jsc_sign > 0) & (voc_sign < 0) or (jsc_sign < 0) & (voc_sign < 0):
                    fwd_j.append(data_fwd[-1, 1])
                    rev_j.append(data_fwd[0, 1])
            elif scan_dir == "rev":
                data_rev = np.genfromtxt(
                    file,
                    delimiter="\t",
                    skip_header=1,
                    skip_footer=NUM_COLS,
                    usecols=(2, 4),
                )
                data_rev = data_rev[~np.isnan(data_rev).any(axis=1)]
                ax1.plot(
                    data_rev[:, 0],
                    data_rev[:, 1],
                    label=pixels[data_index],
                    c=cmap(pixels[data_index] * c_div),
                    lw=2.0,
                )
                if (
                    (jsc_sign > 0) & (voc_sign > 0)
                    or not (jsc_sign > 0) & (voc_sign < 0)
                    and (jsc_sign < 0) & (voc_sign > 0)
                ):
                    fwd_j.append(data_rev[-1, 1])
                    rev_j.append(data_rev[0, 1])
                elif (jsc_sign > 0) & (voc_sign < 0) or (jsc_sign < 0) & (voc_sign < 0):
                    fwd_j.append(data_rev[0, 1])
                    rev_j.append(data_rev[-1, 1])

        # Format the axes
        ax1.tick_params(direction="in", top=True, right=True, labelsize="small")
        ax1.set_xlabel("Applied bias (V)", fontsize="small")
        ax1.set_ylabel("J (mA/cm^2)", fontsize="small")

        # Adjust plot width to add legend outside plot area
        box = ax1.get_position()
        ax1.set_position([box.x0, box.y0, box.width * 0.85, box.height])
        legend_handles, legend_labels = ax1.get_legend_handles_labels()
        lgd = ax1.legend(
            legend_handles,
            legend_labels,
            loc="upper left",
            bbox_to_anchor=(1, 1),
            title="pixel #",
            fontsize="small",
        )

        # Format the figure layout, save to file, and add to ppt
        image_png = os.path.join(image_folder, f"jv_all_{labels[index]}.png")
        image_svg = os.path.join(image_folder, f"jv_all_{labels[index]}.svg")
        fig.savefig(image_png, bbox_extra_artists=(lgd,), bbox_inches="tight")
        fig.savefig(image_svg, bbox_extra_artists=(lgd,), bbox_inches="tight")
        data_slide.shapes.add_picture(
            image_png,
            left=LEFTS[str(index % 4)],
            top=TOPS[str(index % 4)],
            height=IMAGE_HEIGHT,
        )

        # Close figure
        plt.close(fig)


def plot_best_jvs_by_variable_value(best_pixels):
    """Plot JV curves of best pixels by variable value.

    Parameters
    ----------
    best_pixels : pandas.DataFrame
        data frame of best pixels for each variable value
    """
    # create lists of varibales and values for labelling figures
    variables = list(best_pixels["variable"])
    values = list(best_pixels["value"])
    labels = list(best_pixels["label"])
    jsc_signs = list(np.sign(best_pixels["jmpp"]))
    voc_signs = list(np.sign(best_pixels["voc"]))

    # Loop for iterating through best pixels dataframe and picking out JV data
    # files. Each plot contains forward and reverse sweeps, both light and dark.
    for index, (file, scan_dir) in enumerate(
        zip(best_pixels["relativepath"], best_pixels["scandirection"])
    ):
        # Create a new slide after every four graphs are produced
        if index % 4 == 0:
            data_slide = title_image_slide(
                prs, f"Best pixel JVs, page {int(index / 4)}"
            )

        # Create figure, axes, y=0 line, and title
        fig, ax1 = plt.subplots(1, 1, figsize=(A4_WIDTH / 2, A4_HEIGHT / 2), dpi=300)
        ax1.axhline(0, lw=0.5, c="black")
        ax1.axvline(0, lw=0.5, c="black")
        ax1.set_title(
            f"{variables[index]}, {values[index]}, {labels[index]}",
            fontdict={"fontsize": "small"},
        )

        # Import data for each pixel and plot on axes, ignoring errors. If
        # data in a file can't be plotted just ignore it.
        # TODO: handle case with > 2 scans
        if scan_dir == "rev":
            jv_light_rev_path = file
            if file.endswith("liv1"):
                jv_light_fwd_path = file.replace("liv1", "liv2")
            elif file.endswith("liv2"):
                jv_light_fwd_path = file.replace("liv2", "liv1")
        elif scan_dir == "fwd":
            jv_light_fwd_path = file
            if file.endswith("liv1"):
                jv_light_rev_path = file.replace("liv1", "liv2")
            elif file.endswith("liv2"):
                jv_light_rev_path = file.replace("liv2", "liv1")

        with contextlib.suppress(OSError, NameError):
            jv_light_rev_data = np.genfromtxt(
                jv_light_rev_path,
                delimiter="\t",
                skip_header=1,
                skip_footer=NUM_COLS,
                usecols=(2, 4),
            )
            jv_light_fwd_data = np.genfromtxt(
                jv_light_fwd_path,
                delimiter="\t",
                skip_header=1,
                skip_footer=NUM_COLS,
                usecols=(2, 4),
            )
            jv_dark_rev_data = np.genfromtxt(
                jv_light_rev_path.replace("liv", "div"),
                delimiter="\t",
                skip_header=1,
                skip_footer=NUM_COLS,
                usecols=(2, 4),
            )
            jv_dark_fwd_data = np.genfromtxt(
                jv_light_fwd_path.replace("liv", "div"),
                delimiter="\t",
                skip_header=1,
                skip_footer=NUM_COLS,
                usecols=(2, 4),
            )

            jv_light_rev_data = jv_light_rev_data[
                ~np.isnan(jv_light_rev_data).any(axis=1)
            ]
            jv_light_fwd_data = jv_light_fwd_data[
                ~np.isnan(jv_light_fwd_data).any(axis=1)
            ]

            jv_dark_rev_data = jv_dark_rev_data[~np.isnan(jv_dark_rev_data).any(axis=1)]
            jv_dark_fwd_data = jv_dark_fwd_data[~np.isnan(jv_dark_fwd_data).any(axis=1)]

            # plot light J-V curves
            ax1.plot(
                jv_light_rev_data[:, 0],
                jv_light_rev_data[:, 1],
                label="rev",
                c="red",
                lw=2.0,
            )
            ax1.plot(
                jv_light_fwd_data[:, 0],
                jv_light_fwd_data[:, 1],
                label="fwd",
                c="black",
                lw=2.0,
            )

            # find y-limits for plotting
            fwd_j = []
            rev_j = []
            if (jsc_signs[index] > 0) & (voc_signs[index] > 0):
                fwd_j.append(jv_light_rev_data[-1, 1])
                rev_j.append(jv_light_rev_data[0, 1])
                fwd_j.append(jv_light_fwd_data[0, 1])
                rev_j.append(jv_light_fwd_data[-1, 1])
            elif (jsc_signs[index] > 0) & (voc_signs[index] < 0):
                fwd_j.append(jv_light_rev_data[0, 1])
                rev_j.append(jv_light_rev_data[-1, 1])
                fwd_j.append(jv_light_fwd_data[-1, 1])
                rev_j.append(jv_light_fwd_data[0, 1])
            elif (jsc_signs[index] < 0) & (voc_signs[index] > 0):
                fwd_j.append(jv_light_rev_data[-1, 1])
                rev_j.append(jv_light_rev_data[0, 1])
                fwd_j.append(jv_light_fwd_data[0, 1])
                rev_j.append(jv_light_fwd_data[-1, 1])
            elif (jsc_signs[index] < 0) & (voc_signs[index] < 0):
                fwd_j.append(jv_light_rev_data[0, 1])
                rev_j.append(jv_light_rev_data[-1, 1])
                fwd_j.append(jv_light_fwd_data[-1, 1])
                rev_j.append(jv_light_fwd_data[0, 1])

            ax1.plot(
                jv_dark_rev_data[:, 0],
                jv_dark_rev_data[:, 1],
                label="rev",
                c="orange",
                lw=2.0,
            )
            ax1.plot(
                jv_dark_fwd_data[:, 0],
                jv_dark_fwd_data[:, 1],
                label="fwd",
                c="blue",
                lw=2.0,
            )

            # Format the axes
            ax1.tick_params(direction="in", top=True, right=True, labelsize="small")
            ax1.set_xlabel("Applied bias (V)", fontsize="small")
            ax1.set_ylabel("J (mA/cm^2)", fontsize="small")
            ax1.legend(loc="best")

            # Format the figure layout, save to file, and add to ppt
            image_png = os.path.join(
                image_folder, f"jv_best_{variables[index]}_{variables[index]}.png"
            )
            image_svg = os.path.join(
                image_folder, f"jv_best_{variables[index]}_{variables[index]}.svg"
            )
            fig.tight_layout()
            fig.savefig(image_png)
            fig.savefig(image_svg)
            data_slide.shapes.add_picture(
                image_png,
                left=LEFTS[str(index % 4)],
                top=TOPS[str(index % 4)],
                height=IMAGE_HEIGHT,
            )

            # Close figure
            plt.close(fig)


def plot_all_jvs(all_jv_groups):
    """Plot all JV curves.

    Parameters
    ----------
    all_jv_groups : pandas.GroupBy
        all jv scans grouped by pixel number
    """
    # get parameters for plot formatting
    c_div = 1 / 4

    # Create figures, save images and add them to powerpoint slide
    for index, (_, group) in enumerate(all_jv_groups):
        # Create a new slide after every four graphs are produced
        if index % 4 == 0:
            data_slide = title_image_slide(prs, f"All JV scans, page {int(index / 4)}")

        label = group["label"].unique()[0]
        pixel = group["pixel"].unique()[0]
        variable = group["variable"].unique()[0]
        value = group["value"].unique()[0]

        fig, ax1 = plt.subplots(1, 1, figsize=(A4_WIDTH / 2, A4_HEIGHT / 2), dpi=300)
        ax1.axhline(0, lw=0.5, c="black")
        ax1.axvline(0, lw=0.5, c="black")
        ax1.set_title(
            f"{label}, {variable}, {value}, pixel {pixel}",
            fontdict={"fontsize": "small"},
        )

        jsc_signs, jsc_counts = np.unique(np.sign(group["jmpp"]), return_counts=True)
        voc_signs, voc_counts = np.unique(np.sign(group["voc"]), return_counts=True)
        if len(jsc_signs) == 1:
            jsc_sign = jsc_signs[0]
        else:
            max_ix = np.argmax(jsc_counts)
            jsc_sign = jsc_signs[max_ix]
        if len(voc_signs) == 1:
            voc_sign = voc_signs[0]
        else:
            max_ix = np.argmax(voc_counts)
            voc_sign = voc_signs[max_ix]

        # load data for each pixel and plot on axes
        fwd_j = []
        rev_j = []
        for file, scan_dir, scannumber, intensity in zip(
            group["relativepath"],
            group["scandirection"],
            group["scannumber"],
            group["intensity"],
        ):
            if scan_dir == "rev":
                data_rev = np.genfromtxt(
                    file,
                    delimiter="\t",
                    skip_header=1,
                    skip_footer=NUM_COLS,
                    usecols=(2, 4),
                )
                data_rev = data_rev[~np.isnan(data_rev).any(axis=1)]
                if intensity == 0:
                    ax1.plot(
                        data_rev[:, 0],
                        data_rev[:, 1],
                        label=f"{scannumber} {scan_dir} dark",
                        c="black",
                        lw=1.5,
                        ls="--",
                    )
                else:
                    ax1.plot(
                        data_rev[:, 0],
                        data_rev[:, 1],
                        label=f"{scannumber} {scan_dir}",
                        c=cmap(scannumber * c_div),
                        lw=1.5,
                        ls="--",
                    )
                if (
                    (jsc_sign > 0) & (voc_sign > 0)
                    or not (jsc_sign > 0) & (voc_sign < 0)
                    and (jsc_sign < 0) & (voc_sign > 0)
                ):
                    fwd_j.append(data_rev[-1, 1])
                    rev_j.append(data_rev[0, 1])
                elif (jsc_sign > 0) & (voc_sign < 0) or (jsc_sign < 0) & (voc_sign < 0):
                    fwd_j.append(data_rev[0, 1])
                    rev_j.append(data_rev[-1, 1])
            elif scan_dir == "fwd":
                data_fwd = np.genfromtxt(
                    file,
                    delimiter="\t",
                    skip_header=1,
                    skip_footer=NUM_COLS,
                    usecols=(2, 4),
                )
                data_fwd = data_fwd[~np.isnan(data_fwd).any(axis=1)]
                if intensity == 0:
                    ax1.plot(
                        data_fwd[:, 0],
                        data_fwd[:, 1],
                        label=f"{scannumber} {scan_dir} dark",
                        c="black",
                        lw=1.5,
                    )
                else:
                    ax1.plot(
                        data_fwd[:, 0],
                        data_fwd[:, 1],
                        label=f"{scannumber} {scan_dir}",
                        c=cmap(scannumber * c_div),
                        lw=1.5,
                    )
                if (
                    (jsc_sign > 0) & (voc_sign > 0)
                    or not (jsc_sign > 0) & (voc_sign < 0)
                    and (jsc_sign < 0) & (voc_sign > 0)
                ):
                    fwd_j.append(data_fwd[0, 1])
                    rev_j.append(data_fwd[-1, 1])
                elif (jsc_sign > 0) & (voc_sign < 0) or (jsc_sign < 0) & (voc_sign < 0):
                    fwd_j.append(data_fwd[-1, 1])
                    rev_j.append(data_fwd[0, 1])

        # Format the axes
        ax1.tick_params(direction="in", top=True, right=True, labelsize="small")
        ax1.set_xlabel("Applied bias (V)", fontsize="small")
        ax1.set_ylabel("J (mA/cm^2)", fontsize="small")

        # Adjust plot width to add legend outside plot area
        box = ax1.get_position()
        ax1.set_position([box.x0, box.y0, box.width * 0.85, box.height])
        legend_handles, legend_labels = ax1.get_legend_handles_labels()
        lgd = ax1.legend(
            legend_handles,
            legend_labels,
            loc="upper left",
            title="scannumber #",
            bbox_to_anchor=(1, 1),
            fontsize="small",
        )

        # Format the figure layout, save to file, and add to ppt
        image_png = os.path.join(image_folder, f"jv_all_{label}.png")
        image_svg = os.path.join(image_folder, f"jv_all_{label}.svg")
        fig.savefig(image_png, bbox_extra_artists=(lgd,), bbox_inches="tight")
        fig.savefig(image_svg, bbox_extra_artists=(lgd,), bbox_inches="tight")
        data_slide.shapes.add_picture(
            image_png,
            left=LEFTS[str(index % 4)],
            top=TOPS[str(index % 4)],
            height=IMAGE_HEIGHT,
        )

        # Close figure
        plt.close(fig)


# parse args
args = parse()

# create a logger
logger = create_logger(args.folder, args.debug)

FIX_YMIN_0 = True if args.fix_ymin_0 == "yes" else False

# format data if from Python program
(
    analysis_folder,
    experiment_time,
    username,
    experiment_title,
) = format_folder(pathlib.Path(args.folder))

# generate log file
log_filepath = generate_log(analysis_folder)

# Define folder and file paths
folderpath, log_filename = os.path.split(log_filepath)
folderpath_split = recursive_path_split(folderpath)
if folderpath_split[-1] == "LOGS":
    raise ValueError("the log file must be in the same folder as the jv data!")

# change cwd to same folder as log file
os.chdir(folderpath)

# Create folders for storing files generated during analysis
logger.info("Creating analysis folder...")
analysis_folder = os.path.join(folderpath, "Analysis")
image_folder = os.path.join(analysis_folder, "Figures")
if os.path.exists(analysis_folder):
    pass
else:
    os.makedirs(analysis_folder)
if os.path.exists(image_folder):
    pass
else:
    os.makedirs(image_folder)


# Get username, date, and title from folderpath for the ppt title page
exp_date = time.strftime("%A %B %d %Y", time.localtime(experiment_time))

# Set physical constants
kB = scipy.constants.Boltzmann
q = scipy.constants.elementary_charge
c = scipy.constants.speed_of_light
h = scipy.constants.Planck
T = 300

# Read in data from JV log file
logger.info("Loading log file...")
all_data = pd.read_csv(log_filepath, delimiter="\t", header=0)

# format header names so they can be read as attributes
names = []
for name in all_data.columns:
    ix = name.find("(")
    if ix != -1:
        # ignore everything from the first parenthesis
        name = name[:ix]
    # remove all special characters
    name = name.lower().translate(
        {ord(c): "" for c in r"!@#$%^&*()[]{};:,./<>?\|`~-=_+ "}
    )
    names.append(name)
all_data.columns = names

NUM_COLS = len(all_data.columns)


# Create a powerpoint presentation to add figures to.
logger.info("Creating powerpoint file...")
prs = Presentation()

# Add title page with experiment title, date, and username.
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = experiment_title
subtitle.text = f"{exp_date}\n{username}"

# Add blank slide for table of experimental details.
blank_slide_layout = prs.slide_layouts[6]
table_slide = prs.slides.add_slide(blank_slide_layout)
table_shapes = table_slide.shapes
table_rows = len(all_data["label"].unique()) + 1
table_cols = 6
table_left = Inches(0.15)
table_top = Inches(0.02)
table_width = prs.slide_width - Inches(0.25)
table_height = prs.slide_height - Inches(0.05)
table = table_shapes.add_table(
    table_rows, table_cols, table_left, table_top, table_width, table_height
).table

# set column widths
table.columns[0].width = Inches(0.8)
table.columns[1].width = Inches(0.6)
table.columns[2].width = Inches(2.2)
table.columns[3].width = Inches(2.2)
table.columns[4].width = Inches(2.2)
table.columns[5].width = Inches(1.7)

# write column headings
table.cell(0, 0).text = "label"
table.cell(0, 1).text = "sub"
table.cell(0, 2).text = "HTM"
table.cell(0, 3).text = "perovskite"
table.cell(0, 4).text = "ETM"
table.cell(0, 5).text = "metal"

# Define dimensions used for adding images to slides
A4_HEIGHT = 7.5
A4_WIDTH = 10
IMAGE_HEIGHT = prs.slide_height * 0.95 / 2
IMAGE_WIDTH = prs.slide_width * 0.95 / 2

# Create dictionaries that define where to put images on slides in the ppt
LEFTS = {
    "0": Inches(0),
    "1": prs.slide_width - IMAGE_WIDTH,
    "2": Inches(0),
    "3": prs.slide_width - IMAGE_WIDTH,
}
TOPS = {
    "0": prs.slide_height * 0.05,
    "1": prs.slide_height * 0.05,
    "2": prs.slide_height - IMAGE_HEIGHT,
    "3": prs.slide_height - IMAGE_HEIGHT,
}


# Sort data
logger.info("Sorting and filtering data...")

# filter out pixels that yield a non-physical quasi-FF or Jss from analysis
filter_groups = all_data.groupby(["label", "pixel"])
filtered_data_by_ss = filter_groups.filter(
    lambda x: all((x["quasiff"] >= 0) & (x["quasiff"] <= 1) & (x["jss"] < 50))
)

sorted_data = filtered_data_by_ss.sort_values(
    ["label", "pixel", "pce"], ascending=[True, True, False]
)

# Fill in label column of device info table in ppt
table_info = sorted_data.drop_duplicates(["label"])
for table_row, (_, row) in enumerate(table_info.iterrows()):
    table_row += 1
    table.cell(table_row, 0).text = f"{row.label}"
    table.cell(table_row, 1).text = f"{row.substrate}"
    table.cell(table_row, 2).text = f"{row.htm}"
    table.cell(table_row, 3).text = f"{row.perovskite}"
    table.cell(table_row, 4).text = f"{row.etm}"
    table.cell(table_row, 5).text = f"{row.metal}"


# further filter data based on jv scans
filtered_data = sorted_data[
    (sorted_data.intensity > 0)
    & (sorted_data.ff > 0)
    & (sorted_data.ff < 1)
    & (np.absolute(sorted_data.jsc) > 0.01)
    & (not isinstance(sorted_data.jsc, str))
    & (not isinstance(sorted_data.voc, str))
    & (not isinstance(sorted_data.ff, str))
]
filtered_data_fwd = filtered_data[filtered_data.scandirection == "fwd"]
filtered_data_rev = filtered_data[filtered_data.scandirection == "rev"]
filtered_data = filtered_data.drop_duplicates(["label", "pixel", "scandirection"])

filtered_data_fwd = filtered_data_fwd.drop_duplicates(["label", "pixel"])
filtered_data_rev = filtered_data_rev.drop_duplicates(["label", "pixel"])

# Drop pixels only working in one scannumber scandirection.
# First get the inner merge of the label and pixel columns, i.e. drop rows
# where label and pixel combination only occurs in one scannumber scandirection.
filtered_data_fwd_t = filtered_data_fwd[["label", "pixel"]].merge(
    filtered_data_rev[["label", "pixel"]], on=["label", "pixel"], how="inner"
)
filtered_data_rev_t = filtered_data_rev[["label", "pixel"]].merge(
    filtered_data_fwd[["label", "pixel"]], on=["label", "pixel"], how="inner"
)

# Then perform inner merge of full filtered data frames with the merged
# label and pixel dataframes to get back all pixel data that work in both
# scannumber directions
filtered_data_fwd = filtered_data_fwd.merge(
    filtered_data_fwd_t, on=["label", "pixel"], how="inner"
)
filtered_data_rev = filtered_data_rev.merge(
    filtered_data_rev_t, on=["label", "pixel"], how="inner"
)

spo_data = filtered_data_by_ss[
    (np.absolute(sorted_data.vss) > 0) & (np.absolute(sorted_data.jss) > 0)
]
sjsc_data = filtered_data_by_ss[
    (sorted_data.vss == 0)
    & (np.absolute(sorted_data.jss) > 0)
    & (sorted_data.quasiff < 0.9)
    & (sorted_data.quasiff > 0.1)
]
svoc_data = filtered_data_by_ss[
    (np.absolute(sorted_data.vss) > 0) & (sorted_data.jss == 0)
]


logger.info("Plotting spectra...")
# Get spectrum files
spectrum_files = [f for f in os.listdir(folderpath) if f.endswith("spectrum.txt")]
if len(spectrum_files) > 0:
    spectrum_files.sort()
    plot_spectra(spectrum_files)


logger.info("Plotting boxplots and barcharts...")
jv_params = ["jsc", "voc", "ff", "pce", "vmpp", "jmpp", "rsvfwd"]
spo_params = ["pcess", "pcesspcejv"]
sjsc_params = ["jss", "quasiff"]
svoc_params = ["vss"]

# create boxplots for jv and ss parameters grouped by label
boxplot_index = 0
if not svoc_data.empty:
    boxplot_index, data_slide = plot_boxplots(svoc_data, svoc_params, "SSVoc", "label")
    plt.close("all")
if not sjsc_data.empty:
    boxplot_index, data_slide = plot_boxplots(
        sjsc_data,
        sjsc_params,
        "SSJsc",
        "label",
        start_index=boxplot_index,
        data_slide=data_slide,
    )
    plt.close("all")
if not spo_data.empty:
    boxplot_index, data_slide = plot_boxplots(
        spo_data,
        spo_params,
        "SSPO",
        "label",
        start_index=boxplot_index,
        data_slide=data_slide,
    )
    plt.close("all")
if not filtered_data.empty:
    boxplot_index, data_slide = plot_boxplots(filtered_data, jv_params, "J-V", "label")
    plt.close("all")

# create boxplots for jv and spo parameters grouped by variable value
grouped_filtered_data = filtered_data.groupby(["variable"])
grouped_spo_data = spo_data.groupby(["variable"])
grouped_sjsc_data = sjsc_data.groupby(["variable"])
grouped_svoc_data = svoc_data.groupby(["variable"])
for svoc_name, svoc_group in grouped_svoc_data:
    boxplot_index, data_slide = plot_boxplots(
        svoc_group, svoc_params, "SSVoc", "value", svoc_name
    )
    plt.close("all")
for sjsc_name, sjsc_group in grouped_sjsc_data:
    boxplot_index, data_slide = plot_boxplots(
        sjsc_group,
        sjsc_params,
        "SSJsc",
        "value",
        sjsc_name,
        start_index=boxplot_index,
        data_slide=data_slide,
    )
    plt.close("all")
for spo_name, spo_group in grouped_spo_data:
    boxplot_index, data_slide = plot_boxplots(
        spo_group,
        spo_params,
        "SSPO",
        "value",
        spo_name,
        start_index=boxplot_index,
        data_slide=data_slide,
    )
    plt.close("all")
for jv_name, jv_group in grouped_filtered_data:
    boxplot_index, data_slide = plot_boxplots(
        jv_group, jv_params, "J-V", "value", jv_name
    )
    plt.close("all")

# create countplot for yields grouped by label
data_slide = title_image_slide(prs, "Yields, page 0")
plot_countplots(filtered_data, 0, "label", data_slide)

# create countplot for yields grouped by variable value
for yield_index, (yield_name, yield_group) in enumerate(grouped_filtered_data):
    yield_index += 1
    # create new slide if necessary
    if yield_index % 4 == 0:
        data_slide = title_image_slide(prs, f"Yields, page {int(yield_index / 4)}")
    plot_countplots(filtered_data, yield_index, "value", data_slide, yield_name)
    yield_index += 1


# plot steady-state data
logger.info("Plotting steady-state data...")
plot_stabilisation(spo_data, "Steady-state power output", "spo")
plot_stabilisation(sjsc_data, "Steady-state Jsc", "sjsc")
plot_stabilisation(svoc_data, "Steady-state Voc", "svoc")


logger.info("Plotting JV curves...")
# Group data by label and sort ready to plot graph of all pixels per substrate
re_sort_data = filtered_data.sort_values(["label", "pixel"], ascending=[True, True])
grouped_by_label = re_sort_data.groupby("label")

# Create lists of varibales, values, and labels for labelling figures
substrates = re_sort_data.drop_duplicates(["label"])

plot_best_jvs_by_label(grouped_by_label, substrates)

# filter dataframe to leave only the best pixel for each variable value
sort_best_pixels = filtered_data.sort_values(
    ["variable", "value", "pce"], ascending=[True, True, False]
)
best_pixels = sort_best_pixels.drop_duplicates(["variable", "value"])

plot_best_jvs_by_variable_value(best_pixels)


all_jvs_sorted = filtered_data_by_ss.sort_values(
    ["label", "intensity", "pixel", "scannumber"], ascending=[True, False, True, True]
)
all_jvs_groups = all_jvs_sorted[
    (all_jvs_sorted.jsc != 0)
    & (all_jvs_sorted.voc != 0)
    & (all_jvs_sorted.pce != 0)
    & (all_jvs_sorted.ff != 0)
].groupby(["label", "pixel"])

plot_all_jvs(all_jvs_groups)

# Save powerpoint presentation
logger.info("Saving powerpoint presentation...")
prs.save(f"{log_filepath}".replace(".log", "_summary.pptx"))

logger.info("Analysis complete!")
